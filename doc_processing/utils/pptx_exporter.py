"""Utility for exporting PPTX slides as PDFs using LibreOffice."""
import subprocess
import os
import logging
import tempfile
from pathlib import Path
import uuid
from typing import Optional, Union # Import Union and Optional

class PPTXExporter:
    """Exports PPTX slides as PDFs using LibreOffice."""

    def __init__(self, libreoffice_path: str = "/Applications/LibreOffice.app/Contents/MacOS/soffice"):
        """Initialize the PPTX exporter.

        Args:
            libreoffice_path: Path to the LibreOffice executable.
        """
        self.logger = logging.getLogger(__name__)
        self.libreoffice_path = libreoffice_path

        if not Path(self.libreoffice_path).exists():
            self.logger.error(f"LibreOffice executable not found at {self.libreoffice_path}")
            raise FileNotFoundError(f"LibreOffice executable not found at {self.libreoffice_path}")

        self.logger.info(f"Initialized PPTXExporter (libreoffice_path={self.libreoffice_path})")

    def export_to_pdf(self, pptx_path: Union[str, Path]) -> Optional[Path]:
        """Export a PPTX file to PDF.

        Args:
            pptx_path: Path to the input PPTX file.

        Returns:
            Path to the generated PDF file, or None if export failed.
        """
        pptx_path = Path(pptx_path)
        if not pptx_path.exists():
            self.logger.error(f"PPTX file not found at {pptx_path}")
            return None

        # Create a temporary directory for the output PDF
        temp_dir = Path(tempfile.gettempdir())
        output_pdf_path = temp_dir / f"{uuid.uuid4()}.pdf"

        command = [
            self.libreoffice_path,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(temp_dir),
            str(pptx_path)
        ]

        try:
            self.logger.info(f"Executing LibreOffice command: {' '.join(command)}")
            result = subprocess.run(command, capture_output=True, text=True, check=True)
            self.logger.debug(f"LibreOffice stdout:\n{result.stdout}")
            self.logger.debug(f"LibreOffice stderr:\n{result.stderr}")

            if output_pdf_path.exists():
                self.logger.info(f"Successfully exported {pptx_path} to {output_pdf_path}")
                return output_pdf_path
            else:
                self.logger.error(f"LibreOffice command finished but PDF not found at {output_pdf_path}")
                return None

        except FileNotFoundError:
            self.logger.error(f"LibreOffice executable not found at {self.libreoffice_path}. Is LibreOffice installed and the path correct?")
            return None
        except subprocess.CalledProcessError as e:
            self.logger.error(f"LibreOffice command failed with error: {e}")
            self.logger.error(f"Stderr:\n{e.stderr}")
            self.logger.error(f"Stdout:\n{e.stdout}")
            return None
        except Exception as e:
            self.logger.error(f"An unexpected error occurred during PDF export: {e}")
            return None

    def cleanup_pdf(self, pdf_path: Path) -> None:
        """Clean up the temporary PDF file."""
        try:
            if pdf_path.exists():
                pdf_path.unlink()
                self.logger.info(f"Cleaned up temporary PDF file: {pdf_path}")
        except Exception as e:
            self.logger.error(f"Error cleaning up temporary PDF file {pdf_path}: {e}")

if __name__ == '__main__':
    # Example Usage (requires a dummy.pptx file in the current directory)
    exporter = PPTXExporter()
    dummy_pptx = Path("dummy.pptx")

    # Create a dummy pptx file for testing if it doesn't exist
    if not dummy_pptx.exists():
        try:
            from pptx import Presentation
            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = "Dummy Presentation"
            subtitle.text = "For testing PPTXExporter"
            prs.save(dummy_pptx)
            print(f"Created dummy PPTX file: {dummy_pptx}")
        except ImportError:
            print("python-pptx not installed. Cannot create dummy PPTX. Please install it (`poetry add python-pptx`) or provide your own dummy.pptx.")
            exit()
        except Exception as e:
            print(f"Error creating dummy PPTX file: {e}")
            exit()


    pdf_output = exporter.export_to_pdf(dummy_pptx)

    if pdf_output:
        print(f"PDF exported successfully to: {pdf_output}")
        # Clean up the dummy PDF
        # exporter.cleanup_pdf(pdf_output)
        # print(f"Cleaned up {pdf_output}")
    else:
        print("PDF export failed.")

    # Clean up the dummy pptx file
    # try:
    #     if dummy_pptx.exists():
    #         dummy_pptx.unlink()
    #         print(f"Cleaned up dummy PPTX file: {dummy_pptx}")
    # except Exception as e:
    #     print(f"Error cleaning up dummy PPTX file {dummy_pptx}: {e}")