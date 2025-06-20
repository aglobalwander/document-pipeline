"""Script to create a dummy PPTX file for testing."""
from pathlib import Path
from pptx import Presentation
from typing import Union # Import Union

def create_dummy_pptx(output_path: Union[str, Path] = "slides/team_deck.pptx"):
    """Creates a dummy PPTX file with slides and notes."""
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    # Slide 1 with notes
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = "Team Introduction"
    subtitle1.text = "Overview of the team and roles"
    notes_slide1 = slide1.notes_slide
    notes_slide1.notes_text_frame.text = "Notes for slide 1: Introduce team members and their primary responsibilities."

    # Slide 2 with notes
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    title2 = slide2.shapes.title
    body_shape2 = slide2.shapes.placeholders[1]
    tf2 = body_shape2.text_frame
    title2.text = "Project Goals"
    tf2.add_paragraph().text = "Goal 1: Complete feature X"
    tf2.add_paragraph().text = "Goal 2: Improve performance by Y%"
    notes_slide2 = slide2.notes_slide
    notes_slide2.notes_text_frame.text = "Notes for slide 2: Discuss key project objectives and success metrics."

    # Slide 3 without notes
    blank_slide_layout = prs.slide_layouts[6]
    slide3 = prs.slides.add_slide(blank_slide_layout)
    # No notes for this slide

    prs.save(output_path)
    print(f"Created dummy PPTX file: {output_path}")

if __name__ == "__main__":
    try:
        from typing import Union # Import Union here for the script
        create_dummy_pptx()
    except ImportError:
        print("python-pptx not installed. Cannot create dummy PPTX. Please install it (`poetry add python-pptx`).")
    except Exception as e:
        print(f"An error occurred: {e}")