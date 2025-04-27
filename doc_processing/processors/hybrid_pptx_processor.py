"""Hybrid PPTX processor that combines PDF export and speaker notes extraction."""
import logging
import tempfile
from pathlib import Path
from typing import Dict, Any, Optional, List, Union

from doc_processing.embedding.base import BaseProcessor
from doc_processing.utils.pptx_exporter import PPTXExporter
from doc_processing.utils.pptx_notes import PPTXNotesExtractor

class HybridPPTXProcessor(BaseProcessor):
    """Processes PPTX files by exporting to PDF and extracting speaker notes."""

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """Initialize the Hybrid PPTX processor.

        Args:
            config: Configuration options.
                libreoffice_path: Path to the LibreOffice executable (for PDF export).
        """
        super().__init__(config or {})
        self.logger = logging.getLogger(__name__)

        libreoffice_path = self.config.get("libreoffice_path", "/Applications/LibreOffice.app/Contents/MacOS/soffice")
        self.pptx_exporter = PPTXExporter(libreoffice_path=libreoffice_path)
        self.pptx_notes_extractor = PPTXNotesExtractor()

        self.logger.info("Initialized HybridPPTXProcessor")

    def process(self, document: Dict[str, Any]) -> Dict[str, Any]:
        """Process a PPTX document.

        Exports the PPTX to PDF and extracts speaker notes.
        Adds 'pdf_path' and 'speaker_notes' to the document metadata.

        Args:
            document: Document dictionary with 'input_path' pointing to the PPTX file.

        Returns:
            Processed document with updated metadata.
        """
        input_path = document.get('input_path')
        if not input_path:
            self.logger.error("Document does not have an 'input_path'. Cannot process.")
            return document

        pptx_path = Path(input_path)
        if not pptx_path.exists():
            self.logger.error(f"PPTX file not found at {pptx_path}. Cannot process.")
            return document

        # Ensure metadata exists
        if 'metadata' not in document:
            document['metadata'] = {}

        # 1. Export PPTX to PDF
        pdf_path = self.pptx_exporter.export_to_pdf(pptx_path)
        if pdf_path:
            document['metadata']['pdf_path'] = str(pdf_path)
            self.logger.info(f"PDF exported and path added to metadata: {pdf_path}")
        else:
            self.logger.warning(f"Failed to export PPTX to PDF: {pptx_path}")

        # 2. Extract speaker notes
        notes = self.pptx_notes_extractor.extract_notes(pptx_path)
        if notes:
            document['metadata']['speaker_notes'] = notes
            self.logger.info(f"Speaker notes extracted and added to metadata for {pptx_path}")
        else:
            self.logger.warning(f"No speaker notes extracted for {pptx_path}")

        self.logger.info(f"Finished processing PPTX file: {pptx_path}")
        return document

if __name__ == '__main__':
    # Example Usage (requires a dummy.pptx file with notes)
    processor = HybridPPTXProcessor()
    dummy_pptx_with_notes = Path("dummy_with_notes.pptx") # Use the one created by pptx_notes.py example

    # Create a dummy pptx file with notes for testing if it doesn't exist
    if not dummy_pptx_with_notes.exists():
        try:
            from pptx import Presentation
            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = "Dummy Presentation with Notes"
            subtitle.text = "For testing HybridPPTXProcessor"

            # Add notes to the slide
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "These are some speaker notes for the first slide."

            # Add another slide with no notes
            blank_slide_layout = prs.slide_layouts[6]
            slide2 = prs.slides.add_slide(blank_slide_layout)

            prs.save(dummy_pptx_with_notes)
            print(f"Created dummy PPTX file with notes: {dummy_pptx_with_notes}")
        except ImportError:
            print("python-pptx not installed. Cannot create dummy PPTX with notes. Please install it (`poetry add python-pptx`) or provide your own dummy_with_notes.pptx.")
            exit()
        except Exception as e:
            print(f"Error creating dummy PPTX file with notes: {e}")
            exit()


    # Create a dummy document dictionary
    dummy_document = {
        "input_path": str(dummy_pptx_with_notes),
        "content": None, # Content will be derived from PDF and notes
        "metadata": {
            "filename": dummy_pptx_with_notes.name,
            "filetype": "pptx"
        }
    }

    processed_document = processor.process(dummy_document)

    print("\nProcessed Document:")
    print(processed_document)

    # Clean up the dummy files
    # if 'pdf_path' in processed_document.get('metadata', {}):
    #     processor.pptx_exporter.cleanup_pdf(Path(processed_document['metadata']['pdf_path']))
    # try:
    #     if dummy_pptx_with_notes.exists():
    #         dummy_pptx_with_notes.unlink()
    #         print(f"Cleaned up dummy PPTX file: {dummy_pptx_with_notes}")
    # except Exception as e:
    #     print(f"Error cleaning up dummy PPTX file {dummy_pptx_with_notes}: {e}")