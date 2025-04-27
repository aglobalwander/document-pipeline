"""Utility for extracting speaker notes from PPTX files."""
import logging
from pathlib import Path
from typing import Dict, Any, Optional, List, Union # Import Union

try:
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError
except ImportError:
    Presentation = None
    PackageNotFoundError = None
    logging.warning("python-pptx not installed. PPTX notes extraction will not be available.")

class PPTXNotesExtractor:
    """Extracts speaker notes from PPTX slides."""

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """Initialize the PPTX notes extractor."""
        self.logger = logging.getLogger(__name__)
        self.config = config or {}

        if Presentation is None:
            self.logger.error("python-pptx is not installed. Cannot use PPTXNotesExtractor.")
            raise ImportError("python-pptx is not installed. Please install it (`poetry add python-pptx`).")

        self.logger.info("Initialized PPTXNotesExtractor")

    def extract_notes(self, pptx_path: Union[str, Path]) -> List[Dict[str, Any]]:
        """Extract speaker notes from each slide in a PPTX file.

        Args:
            pptx_path: Path to the input PPTX file.

        Returns:
            A list of dictionaries, where each dictionary contains the slide number
            and the extracted notes for that slide.
        """
        pptx_path = Path(pptx_path)
        if not pptx_path.exists():
            self.logger.error(f"PPTX file not found at {pptx_path}")
            return []

        notes_list = []
        try:
            prs = Presentation(pptx_path)
            for i, slide in enumerate(prs.slides):
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text
                    notes_list.append({
                        "slide_number": i + 1,
                        "notes": notes_text.strip()
                    })
                    self.logger.debug(f"Extracted notes for slide {i + 1}")
                else:
                    notes_list.append({
                        "slide_number": i + 1,
                        "notes": ""
                    })
                    self.logger.debug(f"No notes found for slide {i + 1}")

            self.logger.info(f"Extracted notes from {len(prs.slides)} slides in {pptx_path}")
            return notes_list

        except PackageNotFoundError:
            self.logger.error(f"Invalid PPTX file format: {pptx_path}")
            return []
        except Exception as e:
            self.logger.error(f"An error occurred while extracting notes from {pptx_path}: {e}")
            return []

if __name__ == '__main__':
    # Example Usage (requires a dummy.pptx file with notes)
    extractor = PPTXNotesExtractor()
    dummy_pptx = Path("dummy_with_notes.pptx")

    # Create a dummy pptx file with notes for testing if it doesn't exist
    if not dummy_pptx.exists():
        try:
            from pptx import Presentation
            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = "Dummy Presentation with Notes"
            subtitle.text = "For testing PPTXNotesExtractor"

            # Add notes to the slide
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "These are some speaker notes for the first slide."

            # Add another slide with no notes
            blank_slide_layout = prs.slide_layouts[6]
            slide2 = prs.slides.add_slide(blank_slide_layout)

            prs.save(dummy_pptx)
            print(f"Created dummy PPTX file with notes: {dummy_pptx}")
        except ImportError:
            print("python-pptx not installed. Cannot create dummy PPTX with notes. Please install it (`poetry add python-pptx`) or provide your own dummy_with_notes.pptx.")
            exit()
        except Exception as e:
            print(f"Error creating dummy PPTX file with notes: {e}")
            exit()


    notes = extractor.extract_notes(dummy_pptx)

    if notes:
        print("Extracted Notes:")
        for slide_notes in notes:
            print(f"Slide {slide_notes['slide_number']}: {slide_notes['notes']}")
    else:
        print("Failed to extract notes.")

    # Clean up the dummy pptx file
    # try:
    #     if dummy_pptx.exists():
    #         dummy_pptx.unlink()
    #         print(f"Cleaned up dummy PPTX file: {dummy_pptx}")
    # except Exception as e:
    #     print(f"Error cleaning up dummy PPTX file {dummy_pptx}: {e}")